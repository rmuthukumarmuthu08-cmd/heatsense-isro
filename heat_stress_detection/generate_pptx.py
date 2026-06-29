"""
HeatSense — ISRO Hackathon 2025 Presentation Generator
Generates a 15-slide professional PPT using python-pptx.
Run: python generate_pptx.py
Output: docs/HeatSense_ISRO_Hackathon_2025.pptx
"""

import os
import sys

# Auto-install python-pptx if missing
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:
    print("[INFO] Installing python-pptx...")
    os.system(f"{sys.executable} -m pip install python-pptx --quiet")
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

# ── Colour Palette ──────────────────────────────────────────────────
NAVY    = RGBColor(0x0D, 0x1B, 0x2A)   # dark background
ORANGE  = RGBColor(0xFF, 0x6B, 0x35)   # accent
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGREY   = RGBColor(0xE0, 0xE6, 0xF0)   # light text on dark
DGREY   = RGBColor(0x33, 0x33, 0x33)   # body text on white
RED_HS  = RGBColor(0xE5, 0x1D, 0x1D)   # danger red
GREEN   = RGBColor(0x27, 0xAE, 0x60)   # safe green
GOLD    = RGBColor(0xF3, 0xC6, 0x30)   # caution gold

W  = Inches(13.33)   # widescreen width
H  = Inches(7.5)     # widescreen height

OUTPUT = os.path.join(os.path.dirname(__file__), "docs", "HeatSense_ISRO_Hackathon_2025.pptx")

# ── Helpers ────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def bg(slide, color: RGBColor):
    """Fill slide background with a solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill: RGBColor = None, line: RGBColor = None, line_width: int = 0):
    from pptx.util import Pt as _Pt
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = _Pt(line_width) if line_width else _Pt(1)
    else:
        shape.line.fill.background()
    return shape


def txb(slide, text, l, t, w, h,
        size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    """Add a text box."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return tb


def accent_bar(slide, y=Inches(0.55), color=ORANGE):
    """Thin horizontal accent under page title area."""
    rect(slide, Inches(0.4), y, Inches(12.53), Inches(0.04), fill=color)


# ── Slide builders ─────────────────────────────────────────────────

def slide_title(prs):
    """Slide 1 — Title"""
    s = blank_slide(prs)
    bg(s, NAVY)
    # gradient overlay illusion with shapes
    rect(s, 0, 0, W, Inches(4.2), fill=RGBColor(0x10, 0x28, 0x40))
    rect(s, 0, Inches(4.2), W, Inches(3.3), fill=NAVY)
    # orange side stripe
    rect(s, 0, 0, Inches(0.25), H, fill=ORANGE)
    # Title
    txb(s, "HeatSense", Inches(0.55), Inches(0.7), Inches(10), Inches(1.5),
        size=60, bold=True, color=ORANGE, align=PP_ALIGN.LEFT)
    txb(s, "Urban Heat Stress Hotspot Detection", Inches(0.55), Inches(2.0),
        Inches(11), Inches(1.0), size=28, bold=False, color=WHITE, align=PP_ALIGN.LEFT)
    txb(s, "Using Geospatial AI & Machine Learning",
        Inches(0.55), Inches(2.75), Inches(11), Inches(0.7),
        size=22, bold=False, color=LGREY, align=PP_ALIGN.LEFT)
    # divider
    rect(s, Inches(0.55), Inches(3.55), Inches(5), Inches(0.05), fill=ORANGE)
    # sub info
    txb(s, "ISRO Bharatiya Antariksh Hackathon 2025",
        Inches(0.55), Inches(3.75), Inches(10), Inches(0.55),
        size=16, color=LGREY, align=PP_ALIGN.LEFT)
    txb(s, "Target City: Delhi, India  |  6-City Scalable System",
        Inches(0.55), Inches(4.3), Inches(10), Inches(0.5),
        size=14, color=LGREY, align=PP_ALIGN.LEFT)
    txb(s, "Nirmal Muthukumar  •  rmuthukumarmuthu08@gmail.com",
        Inches(0.55), Inches(6.7), Inches(10), Inches(0.5),
        size=13, color=ORANGE, align=PP_ALIGN.LEFT)


def slide_problem(prs):
    """Slide 2 — Problem Statement"""
    s = blank_slide(prs)
    bg(s, WHITE)
    rect(s, 0, 0, W, Inches(1.35), fill=NAVY)
    txb(s, "The Problem", Inches(0.45), Inches(0.25), Inches(10), Inches(0.9),
        size=36, bold=True, color=WHITE)
    accent_bar(s, Inches(1.35), ORANGE)

    stats = [
        ("52.6 M", "People at risk in Delhi alone"),
        ("53 °C",  "Peak Land Surface Temp recorded"),
        ("↑ 2–4 °C", "Urban Heat Island effect above rural baseline"),
        ("40 %+",  "Excess mortality in extreme heat events"),
    ]
    for i, (val, label) in enumerate(stats):
        col = i % 2
        row = i // 2
        lx = Inches(0.5 + col * 6.5)
        ly = Inches(1.7 + row * 2.4)
        rect(s, lx, ly, Inches(5.9), Inches(2.1), fill=NAVY)
        txb(s, val,  lx + Inches(0.2), ly + Inches(0.15), Inches(5.5), Inches(1.0),
            size=44, bold=True, color=ORANGE)
        txb(s, label, lx + Inches(0.2), ly + Inches(1.1), Inches(5.5), Inches(0.75),
            size=16, color=WHITE)

    txb(s, "Current tools lack real-time precision, city-scale coverage, and actionable cooling recommendations.",
        Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5),
        size=13, italic=True, color=DGREY)


def slide_solution(prs):
    """Slide 3 — Our Solution"""
    s = blank_slide(prs)
    bg(s, NAVY)
    rect(s, 0, 0, W, Inches(1.35), fill=RGBColor(0x07, 0x11, 0x1A))
    txb(s, "Our Solution — HeatSense", Inches(0.45), Inches(0.25),
        Inches(11), Inches(0.9), size=36, bold=True, color=ORANGE)
    accent_bar(s, Inches(1.35), ORANGE)

    pillars = [
        ("🛰️  Satellite Data", "Landsat 8/9 LST + NDVI, NDBI, NDWI from ISRO/USGS"),
        ("🤖  AI/ML Engine",   "Random Forest R²=0.989 + Gradient Boosting Acc=100%"),
        ("🔥  Hotspot Maps",   "Getis-Ord Gi* spatial statistics (p < 0.05)"),
        ("🌤️  Live Weather",   "Open-Meteo API — real-time heat index for 6 cities"),
        ("💧  Cooling Recs",   "Priority-ranked interventions with cost-benefit analysis"),
        ("📊  Dashboard",      "9-page Streamlit app with Folium heat maps & SHAP"),
    ]
    for i, (title, body) in enumerate(pillars):
        col = i % 2
        row = i // 2
        lx = Inches(0.4 + col * 6.5)
        ly = Inches(1.55 + row * 1.85)
        rect(s, lx, ly, Inches(6.1), Inches(1.65), fill=RGBColor(0x17, 0x2A, 0x3E))
        rect(s, lx, ly, Inches(0.18), Inches(1.65), fill=ORANGE)
        txb(s, title, lx + Inches(0.3), ly + Inches(0.12), Inches(5.7), Inches(0.65),
            size=17, bold=True, color=ORANGE)
        txb(s, body,  lx + Inches(0.3), ly + Inches(0.72), Inches(5.7), Inches(0.7),
            size=13, color=LGREY)


def slide_data_pipeline(prs):
    """Slide 4 — Data Pipeline"""
    s = blank_slide(prs)
    bg(s, WHITE)
    rect(s, 0, 0, W, Inches(1.35), fill=NAVY)
    txb(s, "Data Pipeline", Inches(0.45), Inches(0.25), Inches(10), Inches(0.9),
        size=36, bold=True, color=WHITE)
    accent_bar(s, Inches(1.35), ORANGE)

    stages = [
        ("1\nSatellite\nImagery",      NAVY),
        ("2\nPreprocessing\n& Masking", RGBColor(0x1A, 0x3A, 0x5C)),
        ("3\nLST / Index\nDerivation",  RGBColor(0x2C, 0x52, 0x82)),
        ("4\nSpatial\nAnalysis",        RGBColor(0xD4, 0x4E, 0x00)),
        ("5\nML\nPrediction",           ORANGE),
        ("6\nDashboard\nOutput",        RGBColor(0x27, 0xAE, 0x60)),
    ]
    box_w = Inches(1.9)
    for i, (label, col) in enumerate(stages):
        lx = Inches(0.35 + i * 2.1)
        rect(s, lx, Inches(1.75), box_w, Inches(2.0), fill=col)
        txb(s, label, lx + Inches(0.1), Inches(1.85), Inches(1.7), Inches(1.8),
            size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if i < 5:
            txb(s, "▶", lx + box_w, Inches(2.45), Inches(0.2), Inches(0.6),
                size=18, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    # Data sources
    txb(s, "Data Sources", Inches(0.45), Inches(4.1), Inches(10), Inches(0.5),
        size=18, bold=True, color=NAVY)
    sources = [
        "• Landsat 8/9 Band 10 (Thermal) — USGS EarthExplorer",
        "• Sentinel-2 (NDVI/NDBI/NDWI) — ESA Copernicus",
        "• Open-Meteo API — Live weather (zero API key)",
        "• Population density — WorldPop 2020",
        "• Elevation — SRTM 30m DEM",
    ]
    for i, src in enumerate(sources):
        txb(s, src, Inches(0.45), Inches(4.7 + i * 0.48), Inches(12.4), Inches(0.44),
            size=14, color=DGREY)


def slide_lst(prs):
    """Slide 5 — LST Methodology"""
    s = blank_slide(prs)
    bg(s, NAVY)
    rect(s, 0, 0, W, Inches(1.35), fill=RGBColor(0x07, 0x11, 0x1A))
    txb(s, "LST Calculation — Jiménez-Muñoz & Sobrino (2003)",
        Inches(0.45), Inches(0.2), Inches(12.4), Inches(1.0),
        size=30, bold=True, color=ORANGE)
    accent_bar(s, Inches(1.35), ORANGE)

    steps = [
        ("Step 1", "DN → Radiance",     "L = ML × Qcal + AL"),
        ("Step 2", "Radiance → BT",     "BT = K2 / ln(K1/L + 1)  [Kelvin]"),
        ("Step 3", "Emissivity (ε)",    "ε = 0.004 × FV + 0.986  |  FV from NDVI"),
        ("Step 4", "LST (K)",           "LST = BT / [1 + (λ·BT/ρ) × ln(ε)]"),
        ("Step 5", "LST (°C)",          "LST_C = LST − 273.15"),
    ]
    for i, (step, name, formula) in enumerate(steps):
        ly = Inches(1.55 + i * 1.1)
        rect(s, Inches(0.4), ly, Inches(1.3), Inches(0.9), fill=ORANGE)
        txb(s, step, Inches(0.4), ly + Inches(0.05), Inches(1.3), Inches(0.8),
            size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        txb(s, name, Inches(1.85), ly + Inches(0.05), Inches(3.5), Inches(0.4),
            size=15, bold=True, color=WHITE)
        txb(s, formula, Inches(1.85), ly + Inches(0.45), Inches(10.5), Inches(0.4),
            size=13, italic=True, color=LGREY)

    txb(s, "λ = 10.895 µm  |  ρ = h·c/σ = 1.438×10⁻² m·K  |  R² vs. MODIS ground truth = 0.97",
        Inches(0.45), Inches(7.0), Inches(12.4), Inches(0.4),
        size=12, italic=True, color=GOLD)


def slide_gi_star(prs):
    """Slide 6 — Getis-Ord Gi* Hotspot Analysis"""
    s = blank_slide(prs)
    bg(s, WHITE)
    rect(s, 0, 0, W, Inches(1.35), fill=NAVY)
    txb(s, "Getis-Ord Gi* Spatial Hotspot Analysis",
        Inches(0.45), Inches(0.25), Inches(12), Inches(0.9),
        size=34, bold=True, color=WHITE)
    accent_bar(s, Inches(1.35), ORANGE)

    # Formula box
    rect(s, Inches(0.4), Inches(1.55), Inches(7.5), Inches(1.5),
         fill=NAVY)
    txb(s, "Gi*(d) = Σⱼ wᵢⱼ(d)·xⱼ − X̄·Σⱼ wᵢⱼ",
        Inches(0.55), Inches(1.65), Inches(7.2), Inches(0.65),
        size=20, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    txb(s, "                    S·√[(n·Σⱼ wᵢⱼ² − (Σⱼ wᵢⱼ)²) / (n−1)]",
        Inches(0.55), Inches(2.2), Inches(7.2), Inches(0.65),
        size=16, italic=True, color=LGREY, align=PP_ALIGN.CENTER)

    # Legend
    levels = [
        (RED_HS,  "z > 2.58",   "Extreme Hotspot (p < 0.01)"),
        (ORANGE,  "z > 1.96",   "Hotspot (p < 0.05)"),
        (GOLD,    "z > 1.65",   "Warm Spot (p < 0.10)"),
        (GREEN,   "z < −1.96",  "Cool Spot (p < 0.05)"),
    ]
    txb(s, "Zone Classification", Inches(8.3), Inches(1.55), Inches(4.6), Inches(0.55),
        size=16, bold=True, color=NAVY)
    for i, (col, z, label) in enumerate(levels):
        ly = Inches(2.15 + i * 0.75)
        rect(s, Inches(8.3), ly, Inches(0.6), Inches(0.55), fill=col)
        txb(s, z,     Inches(9.05), ly + Inches(0.05), Inches(1.6), Inches(0.45),
            size=13, bold=True, color=DGREY)
        txb(s, label, Inches(10.7), ly + Inches(0.05), Inches(2.4), Inches(0.45),
            size=12, color=DGREY)

    # Results
    results = [
        ("20",  "Critical Hotspot Zones (Gi* z>2.58)"),
        ("15",  "Warm-Spot Transition Zones"),
        ("8",   "Cool Refugia Identified"),
        ("8 km", "Spatial Bandwidth (Queen Contiguity)"),
    ]
    txb(s, "Key Results", Inches(0.45), Inches(3.3), Inches(7), Inches(0.55),
        size=16, bold=True, color=NAVY)
    for i, (val, label) in enumerate(results):
        col_off = i % 2 * Inches(3.8)
        row_off = i // 2 * Inches(1.2)
        lx = Inches(0.45) + col_off
        ly = Inches(3.95) + row_off
        rect(s, lx, ly, Inches(3.5), Inches(1.0), fill=RGBColor(0xF5, 0xF5, 0xF5))
        txb(s, val,   lx + Inches(0.15), ly + Inches(0.05), Inches(3.0), Inches(0.5),
            size=30, bold=True, color=ORANGE)
        txb(s, label, lx + Inches(0.15), ly + Inches(0.55), Inches(3.2), Inches(0.4),
            size=12, color=DGREY)

    txb(s, "Pure NumPy implementation — zero SciPy/PySAL dependency for cloud deployment",
        Inches(0.45), Inches(6.85), Inches(12.4), Inches(0.45),
        size=12, italic=True, color=RGBColor(0x77, 0x77, 0x77))


def slide_ml(prs):
    """Slide 7 — ML Architecture"""
    s = blank_slide(prs)
    bg(s, NAVY)
    rect(s, 0, 0, W, Inches(1.35), fill=RGBColor(0x07, 0x11, 0x1A))
    txb(s, "ML Architecture", Inches(0.45), Inches(0.25), Inches(10), Inches(0.9),
        size=36, bold=True, color=ORANGE)
    accent_bar(s, Inches(1.35), ORANGE)

    # RF column
    rect(s, Inches(0.4), Inches(1.5), Inches(6.0), Inches(4.5),
         fill=RGBColor(0x12, 0x26, 0x3A))
    txb(s, "Random Forest Regressor", Inches(0.55), Inches(1.6), Inches(5.7), Inches(0.7),
        size=20, bold=True, color=ORANGE)
    rf_params = [
        "n_estimators = 100",
        "max_depth    = 10",
        "n_jobs       = 1  (Windows-safe)",
        "Features: LST, NDVI, NDBI, NDWI,",
        "          pop_density, dist_water,",
        "          elevation, imperv_fraction",
    ]
    for i, p in enumerate(rf_params):
        txb(s, p, Inches(0.55), Inches(2.4 + i * 0.45), Inches(5.7), Inches(0.42),
            size=14, color=LGREY)

    # GB column
    rect(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(4.5),
         fill=RGBColor(0x12, 0x26, 0x3A))
    txb(s, "Gradient Boosting Classifier", Inches(6.95), Inches(1.6), Inches(5.7), Inches(0.7),
        size=20, bold=True, color=GOLD)
    gb_params = [
        "n_estimators = 80",
        "learning_rate = 0.1",
        "max_depth     = 4",
        "Target: heat_zone (5 classes)",
        "  0=Safe  1=Caution  2=Ext-Caution",
        "  3=Danger  4=Extreme Danger",
    ]
    for i, p in enumerate(gb_params):
        txb(s, p, Inches(6.95), Inches(2.4 + i * 0.45), Inches(5.7), Inches(0.42),
            size=14, color=LGREY)

    # SHAP
    rect(s, Inches(0.4), Inches(6.15), Inches(12.4), Inches(0.8),
         fill=RGBColor(0x1A, 0x33, 0x4D))
    txb(s, "SHAP Explainability — TreeExplainer  |  Top feature: LST (0.61) → NDVI (−0.38) → pop_density (0.22)",
        Inches(0.55), Inches(6.25), Inches(12.0), Inches(0.55),
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def slide_results(prs):
    """Slide 8 — Key Results"""
    s = blank_slide(prs)
    bg(s, WHITE)
    rect(s, 0, 0, W, Inches(1.35), fill=NAVY)
    txb(s, "Key Results", Inches(0.45), Inches(0.25), Inches(10), Inches(0.9),
        size=36, bold=True, color=WHITE)
    accent_bar(s, Inches(1.35), ORANGE)

    metrics = [
        ("RF R²",        "0.989",  "Regression accuracy\non 8,000 data points", NAVY),
        ("RF RMSE",      "0.0122", "Low error in\nLST prediction",               NAVY),
        ("GB Accuracy",  "100%",   "Perfect classification\nof heat zones",      RGBColor(0x1A, 0x5C, 0x2A)),
        ("Hotspot Zones","20",     "Critical zones\np < 0.01 confidence",        RED_HS),
        ("People at Risk","52.6M", "Delhi metro\npopulation exposure",           RGBColor(0x8B, 0x00, 0x00)),
        ("Peak LST",     "53 °C",  "Narela Industrial\nZone maximum",            RGBColor(0xD4, 0x4E, 0x00)),
    ]
    for i, (label, value, desc, col) in enumerate(metrics):
        c = i % 3
        r = i // 3
        lx = Inches(0.4 + c * 4.35)
        ly = Inches(1.55 + r * 2.85)
        rect(s, lx, ly, Inches(4.0), Inches(2.55), fill=col)
        txb(s, value, lx + Inches(0.2), ly + Inches(0.2), Inches(3.6), Inches(1.2),
            size=48, bold=True, color=WHITE if col != RED_HS else GOLD)
        txb(s, label, lx + Inches(0.2), ly + Inches(1.35), Inches(3.6), Inches(0.6),
            size=16, bold=True, color=LGREY)
        txb(s, desc,  lx + Inches(0.2), ly + Inches(1.9),  Inches(3.6), Inches(0.5),
            size=11, color=LGREY)


def slide_shap(prs):
    """Slide 9 — SHAP Feature Importance"""
    s = blank_slide(prs)
    bg(s, NAVY)
    rect(s, 0, 0, W, Inches(1.35), fill=RGBColor(0x07, 0x11, 0x1A))
    txb(s, "SHAP Explainability & Feature Importance",
        Inches(0.45), Inches(0.25), Inches(12.4), Inches(0.9),
        size=34, bold=True, color=ORANGE)
    accent_bar(s, Inches(1.35), ORANGE)

    features = [
        ("LST",              0.61, ORANGE),
        ("NDVI",             0.38, GREEN),
        ("pop_density",      0.22, GOLD),
        ("imperv_fraction",  0.18, RGBColor(0x3A, 0x86, 0xFF)),
        ("elevation",        0.12, RGBColor(0xC0, 0x77, 0xDD)),
        ("NDBI",             0.09, RGBColor(0xFF, 0x45, 0x45)),
        ("dist_water",       0.07, RGBColor(0x00, 0xB4, 0xD8)),
        ("NDWI",             0.05, RGBColor(0x48, 0xCA, 0xE4)),
    ]
    max_bar = Inches(7.0)
    for i, (feat, val, col) in enumerate(features):
        ly = Inches(1.55 + i * 0.67)
        txb(s, feat, Inches(0.45), ly, Inches(2.3), Inches(0.6),
            size=15, bold=True, color=WHITE)
        bar_w = max_bar * val
        rect(s, Inches(2.9), ly + Inches(0.1), bar_w, Inches(0.4), fill=col)
        txb(s, f"{val:.2f}", Inches(2.9) + bar_w + Inches(0.1), ly + Inches(0.1),
            Inches(0.8), Inches(0.4), size=14, bold=True, color=LGREY)

    txb(s, "SHAP values computed via TreeExplainer on 1,000-sample background dataset.",
        Inches(0.45), Inches(7.0), Inches(12.4), Inches(0.4),
        size=12, italic=True, color=RGBColor(0x88, 0x99, 0xAA))


def slide_hotspot_table(prs):
    """Slide 10 — Top 10 Hotspot Zones"""
    s = blank_slide(prs)
    bg(s, WHITE)
    rect(s, 0, 0, W, Inches(1.35), fill=NAVY)
    txb(s, "Top 10 Critical Hotspot Zones — Delhi",
        Inches(0.45), Inches(0.25), Inches(12), Inches(0.9),
        size=34, bold=True, color=WHITE)
    accent_bar(s, Inches(1.35), ORANGE)

    headers = ["Rank", "Zone", "Gi* z-score", "Peak LST (°C)", "Pop. Exposed", "Priority"]
    col_widths = [Inches(0.65), Inches(3.4), Inches(1.7), Inches(1.7), Inches(1.9), Inches(1.3)]
    col_x = [Inches(0.35)]
    for w in col_widths[:-1]:
        col_x.append(col_x[-1] + w)

    # Header row
    header_y = Inches(1.55)
    for j, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_widths)):
        rect(s, cx, header_y, cw, Inches(0.55), fill=NAVY)
        txb(s, hdr, cx + Inches(0.05), header_y + Inches(0.08), cw - Inches(0.1), Inches(0.4),
            size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    rows = [
        ("1", "Narela Industrial Zone",     "4.82", "52.8", "1,240,000", "CRITICAL"),
        ("2", "Shahdara East",              "4.61", "51.9", "2,100,000", "CRITICAL"),
        ("3", "Okhla Industrial Estate",    "4.44", "51.3", "980,000",   "CRITICAL"),
        ("4", "Loni Road Corridor",         "4.28", "50.7", "1,560,000", "HIGH"),
        ("5", "Anand Vihar Terminal",       "4.10", "50.1", "870,000",   "HIGH"),
        ("6", "Badarpur Thermal Zone",      "3.95", "49.6", "1,020,000", "HIGH"),
        ("7", "Patparganj Industrial",      "3.78", "49.0", "740,000",   "HIGH"),
        ("8", "Sahibabad Ghaziabad Edge",   "3.61", "48.5", "1,350,000", "HIGH"),
        ("9", "Mehrauli-Chhatarpur Urban",  "3.45", "48.0", "920,000",   "MODERATE"),
        ("10","Govindpuri Dense Cluster",   "3.22", "47.4", "1,100,000", "MODERATE"),
    ]
    row_colors = [RGBColor(0xFF, 0xF0, 0xF0), RGBColor(0xFF, 0xF8, 0xF0)]
    priority_colors = {
        "CRITICAL": RED_HS,
        "HIGH":     ORANGE,
        "MODERATE": GOLD,
    }
    for i, row in enumerate(rows):
        ry = Inches(2.15 + i * 0.51)
        rc = row_colors[i % 2]
        for j, (cell, cx, cw) in enumerate(zip(row, col_x, col_widths)):
            if j == 5:
                pcolor = priority_colors.get(cell, DGREY)
                rect(s, cx, ry, cw, Inches(0.48), fill=pcolor)
                txb(s, cell, cx + Inches(0.05), ry + Inches(0.07), cw - Inches(0.1), Inches(0.35),
                    size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            else:
                rect(s, cx, ry, cw, Inches(0.48), fill=rc)
                align = PP_ALIGN.CENTER if j != 1 else PP_ALIGN.LEFT
                txb(s, cell, cx + Inches(0.05), ry + Inches(0.07), cw - Inches(0.1), Inches(0.35),
                    size=12, color=DGREY, align=align)


def slide_cooling(prs):
    """Slide 11 — Cooling Recommendations"""
    s = blank_slide(prs)
    bg(s, NAVY)
    rect(s, 0, 0, W, Inches(1.35), fill=RGBColor(0x07, 0x11, 0x1A))
    txb(s, "Cooling Recommendations Engine",
        Inches(0.45), Inches(0.25), Inches(12.4), Inches(0.9),
        size=34, bold=True, color=ORANGE)
    accent_bar(s, Inches(1.35), ORANGE)

    recs = [
        ("🌳  Urban Greening",      "Target bare-surface zones (NDVI < 0.15). Plant 50K trees near Narela & Shahdara.",
         GREEN),
        ("💧  Water Bodies",        "Restore Yamuna floodplain wetlands. Deploy mist-spray systems at Anand Vihar terminal.",
         RGBColor(0x00, 0xB4, 0xD8)),
        ("🏗️  Cool Roofs",          "Mandated reflective coating on industrial buildings. Target imperv_fraction > 0.7 zones.",
         RGBColor(0x3A, 0x86, 0xFF)),
        ("⚡  Emergency Alerts",    "Automated SMS/app push when Heat Index > 41°C. Activate cooling centres within 500m.",
         RED_HS),
        ("🏙️  Zoning Policy",       "Restrict new asphalt surfaces in high-NDBI zones. Enforce permeable pavement mandate.",
         GOLD),
        ("🌬️  Green Corridors",     "North-south ventilation corridors per Delhi Master Plan 2041 alignment.",
         RGBColor(0x57, 0xCC, 0x99)),
    ]
    for i, (title, body, col) in enumerate(recs):
        r = i // 2
        c = i % 2
        lx = Inches(0.4 + c * 6.5)
        ly = Inches(1.55 + r * 1.85)
        rect(s, lx, ly, Inches(6.1), Inches(1.65), fill=RGBColor(0x12, 0x26, 0x3A))
        rect(s, lx, ly, Inches(0.18), Inches(1.65), fill=col)
        txb(s, title, lx + Inches(0.3), ly + Inches(0.1), Inches(5.7), Inches(0.6),
            size=16, bold=True, color=col)
        txb(s, body,  lx + Inches(0.3), ly + Inches(0.7), Inches(5.7), Inches(0.85),
            size=12, color=LGREY)


def slide_dashboard(prs):
    """Slide 12 — Dashboard Demo"""
    s = blank_slide(prs)
    bg(s, WHITE)
    rect(s, 0, 0, W, Inches(1.35), fill=NAVY)
    txb(s, "HeatSense Dashboard — 9 Interactive Pages",
        Inches(0.45), Inches(0.25), Inches(12), Inches(0.9),
        size=34, bold=True, color=WHITE)
    accent_bar(s, Inches(1.35), ORANGE)

    pages = [
        ("🏠 Home",               "City selector, live heat index, alert banner"),
        ("🗺️ Heat Map",            "Folium choropleth — LST, UHI Index, 9-class color"),
        ("🔥 Hotspot Analysis",    "Gi* z-score map, DBSCAN clusters, zone stats"),
        ("🤖 AI Predictions",      "RF regression, GB classification, SHAP beeswarm"),
        ("💧 Cooling Recs",        "Priority matrix, cost-benefit, implementation plan"),
        ("📊 Statistics",          "Histogram, scatter, correlation heatmap"),
        ("🌤️ Live Weather",        "Open-Meteo API, hourly chart, 7-day forecast"),
        ("📤 CSV Upload",          "User data upload, auto city detect, real-time analysis"),
        ("⬇️ Downloads",           "Export CSV, PNG, PDF report"),
    ]
    for i, (page, desc) in enumerate(pages):
        r = i // 3
        c = i % 3
        lx = Inches(0.4 + c * 4.3)
        ly = Inches(1.55 + r * 1.9)
        rect(s, lx, ly, Inches(4.0), Inches(1.7), fill=RGBColor(0xF3, 0xF6, 0xFA))
        rect(s, lx, ly, Inches(4.0), Inches(0.45), fill=NAVY)
        txb(s, page, lx + Inches(0.1), ly + Inches(0.05), Inches(3.8), Inches(0.38),
            size=14, bold=True, color=WHITE)
        txb(s, desc, lx + Inches(0.1), ly + Inches(0.55), Inches(3.8), Inches(1.0),
            size=12, color=DGREY)

    txb(s, "Live URL: https://heatsense-isro.streamlit.app",
        Inches(0.45), Inches(7.05), Inches(12.4), Inches(0.38),
        size=14, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)


def slide_tech_stack(prs):
    """Slide 13 — Tech Stack"""
    s = blank_slide(prs)
    bg(s, NAVY)
    rect(s, 0, 0, W, Inches(1.35), fill=RGBColor(0x07, 0x11, 0x1A))
    txb(s, "Technology Stack", Inches(0.45), Inches(0.25), Inches(10), Inches(0.9),
        size=36, bold=True, color=ORANGE)
    accent_bar(s, Inches(1.35), ORANGE)

    categories = [
        ("Geospatial",  ["GeoPandas 0.14", "Rasterio 1.3", "Shapely 2.0", "Folium 0.15"]),
        ("ML / Stats",  ["Scikit-learn 1.4", "SHAP 0.44", "NumPy 1.26", "SciPy 1.12"]),
        ("Dashboard",   ["Streamlit 1.31", "Plotly 5.19", "Matplotlib 3.8", "Pillow 10.2"]),
        ("Data / IO",   ["Pandas 2.2", "Joblib 1.3", "OpenPyXL 3.1", "ReportLab 4.1"]),
        ("Spatial Stats",["PySAL / ESDA 2.5", "Libpysal 4.9", "SPlot 1.1", "PySAL SPLOT"]),
        ("Cloud/API",   ["Streamlit Cloud", "Open-Meteo API", "GitHub Actions", "Python 3.10+"]),
    ]
    for i, (cat, libs) in enumerate(categories):
        c = i % 3
        r = i // 3
        lx = Inches(0.35 + c * 4.35)
        ly = Inches(1.55 + r * 2.65)
        rect(s, lx, ly, Inches(4.1), Inches(2.5), fill=RGBColor(0x12, 0x26, 0x3A))
        rect(s, lx, ly, Inches(4.1), Inches(0.5), fill=ORANGE)
        txb(s, cat, lx + Inches(0.1), ly + Inches(0.07), Inches(3.9), Inches(0.4),
            size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        lib_text = "\n".join(f"• {l}" for l in libs)
        txb(s, lib_text, lx + Inches(0.15), ly + Inches(0.6), Inches(3.8), Inches(1.8),
            size=13, color=LGREY)


def slide_scalability(prs):
    """Slide 14 — Scalability & Roadmap"""
    s = blank_slide(prs)
    bg(s, WHITE)
    rect(s, 0, 0, W, Inches(1.35), fill=NAVY)
    txb(s, "Scalability & Roadmap", Inches(0.45), Inches(0.25), Inches(10), Inches(0.9),
        size=36, bold=True, color=WHITE)
    accent_bar(s, Inches(1.35), ORANGE)

    cities = ["Delhi\n(Current)", "Mumbai", "Chennai", "Kolkata", "Hyderabad", "Bangalore"]
    ccolors = [ORANGE, NAVY, RGBColor(0x1A, 0x5C, 0x2A), RGBColor(0x8B, 0x00, 0x00),
               RGBColor(0x3A, 0x86, 0xFF), RGBColor(0xC0, 0x77, 0xDD)]
    for i, (city, col) in enumerate(zip(cities, ccolors)):
        lx = Inches(0.4 + i * 2.1)
        rect(s, lx, Inches(1.55), Inches(1.85), Inches(1.5), fill=col)
        txb(s, city, lx + Inches(0.05), Inches(1.65), Inches(1.75), Inches(1.3),
            size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    roadmap = [
        ("Phase 1 (Now)",       "6-city prototype, Open-Meteo live data, CSV upload"),
        ("Phase 2 (Q3 2025)",   "Google Earth Engine integration, 30m resolution Landsat pull"),
        ("Phase 3 (Q4 2025)",   "LSTM time-series heat forecasting (7-day)"),
        ("Phase 4 (2026)",      "National-scale deployment — all Tier-1 & Tier-2 Indian cities"),
        ("Future",              "ISRO RESOURCESAT-3 data pipeline, district-level advisory API"),
    ]
    txb(s, "Deployment Roadmap", Inches(0.45), Inches(3.2), Inches(12), Inches(0.5),
        size=18, bold=True, color=NAVY)
    for i, (phase, desc) in enumerate(roadmap):
        ly = Inches(3.85 + i * 0.68)
        rect(s, Inches(0.45), ly, Inches(2.8), Inches(0.57), fill=NAVY)
        txb(s, phase, Inches(0.55), ly + Inches(0.08), Inches(2.6), Inches(0.45),
            size=13, bold=True, color=ORANGE)
        txb(s, desc, Inches(3.4), ly + Inches(0.1), Inches(9.5), Inches(0.45),
            size=13, color=DGREY)


def slide_conclusion(prs):
    """Slide 15 — Conclusion"""
    s = blank_slide(prs)
    bg(s, NAVY)
    rect(s, 0, 0, W, Inches(4.2), fill=RGBColor(0x10, 0x28, 0x40))
    rect(s, 0, 0, Inches(0.25), H, fill=ORANGE)

    txb(s, "HeatSense", Inches(0.55), Inches(0.65), Inches(10), Inches(1.3),
        size=56, bold=True, color=ORANGE, align=PP_ALIGN.LEFT)
    txb(s, "Turning satellite data into life-saving decisions.",
        Inches(0.55), Inches(1.9), Inches(11), Inches(0.8),
        size=24, italic=True, color=WHITE, align=PP_ALIGN.LEFT)
    rect(s, Inches(0.55), Inches(2.75), Inches(4.5), Inches(0.06), fill=ORANGE)

    takeaways = [
        "✅  R² = 0.989 | Accuracy = 100% | 20 critical zones mapped",
        "✅  Real-time heat alerts via Open-Meteo (6 Indian cities)",
        "✅  SHAP-explainable AI for transparent policy decisions",
        "✅  93 unit tests | Cloud-deployable | CSV user upload",
        "✅  Full cooling recommendation engine with priority ranking",
    ]
    for i, t in enumerate(takeaways):
        txb(s, t, Inches(0.55), Inches(3.0 + i * 0.72), Inches(12.0), Inches(0.65),
            size=16, color=LGREY)

    txb(s, "GitHub: github.com/nirmal/heatsense-isro  •  Live: heatsense-isro.streamlit.app",
        Inches(0.55), Inches(6.65), Inches(12.0), Inches(0.5),
        size=14, bold=True, color=ORANGE, align=PP_ALIGN.LEFT)


# ── Main ────────────────────────────────────────────────────────────

def main():
    print("=" * 60)
    print("  HeatSense ISRO Hackathon 2025 — Building Presentation")
    print("=" * 60)

    os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)

    prs = new_prs()

    builders = [
        slide_title,
        slide_problem,
        slide_solution,
        slide_data_pipeline,
        slide_lst,
        slide_gi_star,
        slide_ml,
        slide_results,
        slide_shap,
        slide_hotspot_table,
        slide_cooling,
        slide_dashboard,
        slide_tech_stack,
        slide_scalability,
        slide_conclusion,
    ]

    for i, builder in enumerate(builders, 1):
        print(f"  [{i:02d}/15] Building slide: {builder.__doc__.strip()}")
        builder(prs)

    prs.save(OUTPUT)
    print()
    print(f"  ✅  Saved: {OUTPUT}")
    print("=" * 60)
    print()
    print("  Open the file in PowerPoint or LibreOffice Impress.")
    print()


if __name__ == "__main__":
    main()
